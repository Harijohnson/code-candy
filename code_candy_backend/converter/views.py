from django.core.files.storage import default_storage
from django.http import JsonResponse
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
import os
from io import BytesIO
from PIL import Image
import pdfkit
import pandas as pd
import markdown2
from docx2pdf import convert as docx2pdf_convert
from pdf2docx import Converter as PDF2DOCX
from pptx import Presentation
from reportlab.pdfgen import canvas
import fitz  # PyMuPDF
from django.conf import settings


@api_view(['POST'])
def convert_file(request):
    if 'file' not in request.FILES or 'operation' not in request.data:
        return Response({'error': 'File and operation are required'}, status=status.HTTP_400_BAD_REQUEST)

    file = request.FILES['file']
    operation = request.data.get('operation')
    if isinstance(operation, list):
        operation = operation[0]  # Extract the first item if it's a list
    
    file_path = default_storage.save(file.name, file)
    file_path = default_storage.path(file_path)
    file_path = os.path.join(settings.MEDIA_ROOT, file_path)
    try:
        if operation == 'pdf_to_word':
            result = convert_pdf_to_word(file_path)
        elif operation == 'image_to_pdf':
            result = convert_image_to_pdf(file_path)
        elif operation == 'word_to_pdf':
            result = convert_word_to_pdf(file_path)
        elif operation == 'jpg_to_png':
            result = convert_jpg_to_png(file_path)
        elif operation == 'png_to_jpg':
            result = convert_png_to_jpg(file_path)
        elif operation == 'excel_to_csv':
            result = convert_excel_to_csv(file_path)
        elif operation == 'csv_to_excel':
            result = convert_csv_to_excel(file_path)
        elif operation == 'html_to_pdf':
            result = convert_html_to_pdf(file_path)
        elif operation == 'pdf_to_html':
            result = convert_pdf_to_html(file_path)
        elif operation == 'ppt_to_pdf':
            result = convert_ppt_to_pdf(file_path)
        elif operation == 'pdf_to_ppt':
            result = convert_pdf_to_ppt(file_path)
        elif operation == 'text_to_pdf':
            result = convert_text_to_pdf(file_path)
        elif operation == 'pdf_to_text':
            result = convert_pdf_to_text(file_path)
        elif operation == 'tiff_to_jpg':
            result = convert_tiff_to_jpg(file_path)
        elif operation == 'jpg_to_tiff':
            result = convert_jpg_to_tiff(file_path)
        elif operation == 'pdf_to_docx':
            result = convert_pdf_to_docx(file_path)
        elif operation == 'docx_to_pdf':
            result = convert_docx_to_pdf(file_path)
        elif operation == 'markdown_to_pdf':
            result = convert_markdown_to_pdf(file_path)
        elif operation == 'png_to_svg':
            result = convert_png_to_svg(file_path)
        else:
            return Response({'error': 'Invalid conversion type'}, status=status.HTTP_400_BAD_REQUEST)

        return Response({'file_url': result}, status=status.HTTP_200_OK)

    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    finally:
        # Clean up the temporary files
        if os.path.exists(file_path):
            os.remove(file_path)

def convert_pdf_to_word(file_path):
    print('file_path', file_path)
    try:
        output_path = file_path.replace('.pdf', '.docx')
        cv = PDF2DOCX(file_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return default_storage.url(output_path)
    except Exception as e:
        print(f'Exception while converting PDF to Word: {e}')
        raise RuntimeError(f'PDF to Word conversion failed: {e}')

def convert_image_to_pdf(file_path):
    try:
        output_path = file_path.replace('.jpg', '.pdf')
        image = Image.open(file_path)
        image.convert('RGB').save(output_path, 'PDF')
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'Image to PDF conversion failed: {e}')

def convert_word_to_pdf(file_path):
    try:
        output_path = file_path.replace('.docx', '.pdf')
        docx2pdf_convert(file_path, output_path)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'Word to PDF conversion failed: {e}')

def convert_jpg_to_png(file_path):
    try:
        output_path = file_path.replace('.jpg', '.png')
        image = Image.open(file_path)
        image.save(output_path, 'PNG')
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'JPG to PNG conversion failed: {e}')

def convert_png_to_jpg(file_path):
    try:
        output_path = file_path.replace('.png', '.jpg')
        image = Image.open(file_path)
        image.convert('RGB').save(output_path, 'JPEG')
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PNG to JPG conversion failed: {e}')

def convert_excel_to_csv(file_path):
    try:
        output_path = file_path.replace('.xlsx', '.csv')
        df = pd.read_excel(file_path)
        df.to_csv(output_path, index=False)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'Excel to CSV conversion failed: {e}')

def convert_csv_to_excel(file_path):
    try:
        output_path = file_path.replace('.csv', '.xlsx')
        df = pd.read_csv(file_path)
        df.to_excel(output_path, index=False)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'CSV to Excel conversion failed: {e}')

def convert_html_to_pdf(file_path):
    try:
        output_path = file_path.replace('.html', '.pdf')
        pdfkit.from_file(file_path, output_path)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'HTML to PDF conversion failed: {e}')

def convert_pdf_to_html(file_path):
    try:
        output_path = file_path.replace('.pdf', '.html')
        with open(output_path, 'w') as file:
            file.write("<html><body>PDF to HTML conversion is not yet implemented.</body></html>")
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PDF to HTML conversion failed: {e}')

def convert_ppt_to_pdf(file_path):
    try:
        output_path = file_path.replace('.pptx', '.pdf')
        prs = Presentation(file_path)
        c = canvas.Canvas(output_path)
        for slide in prs.slides:
            c.drawString(100, 750, "PowerPoint slide content")  # Placeholder
        c.save()
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PPT to PDF conversion failed: {e}')

def convert_pdf_to_ppt(file_path):
    try:
        output_path = file_path.replace('.pdf', '.pptx')
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "PDF to PowerPoint conversion is not yet implemented."
        prs.save(output_path)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PDF to PPT conversion failed: {e}')

def convert_text_to_pdf(file_path):
    try:
        output_path = file_path.replace('.txt', '.pdf')
        with open(file_path, 'r') as file:
            text = file.read()
        c = canvas.Canvas(output_path)
        c.drawString(100, 750, text)
        c.save()
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'Text to PDF conversion failed: {e}')

def convert_pdf_to_text(file_path):
    try:
        output_path = file_path.replace('.pdf', '.txt')
        pdf_document = fitz.open(file_path)
        with open(output_path, 'w') as file:
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                file.write(page.get_text())
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PDF to Text conversion failed: {e}')

def convert_tiff_to_jpg(file_path):
    try:
        output_path = file_path.replace('.tiff', '.jpg')
        image = Image.open(file_path)
        image.convert('RGB').save(output_path, 'JPEG')
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'TIFF to JPG conversion failed: {e}')

def convert_jpg_to_tiff(file_path):
    try:
        output_path = file_path.replace('.jpg', '.tiff')
        image = Image.open(file_path)
        image.save(output_path, 'TIFF')
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'JPG to TIFF conversion failed: {e}')

def convert_pdf_to_docx(file_path):
    try:
        output_path = file_path.replace('.pdf', '.docx')
        cv = PDF2DOCX(file_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PDF to DOCX conversion failed: {e}')

def convert_docx_to_pdf(file_path):
    try:
        output_path = file_path.replace('.docx', '.pdf')
        docx2pdf_convert(file_path, output_path)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'DOCX to PDF conversion failed: {e}')

def convert_markdown_to_pdf(file_path):
    try:
        output_path = file_path.replace('.md', '.pdf')
        with open(file_path, 'r') as file:
            markdown_text = file.read()
        html_text = markdown2.markdown(markdown_text)
        pdfkit.from_string(html_text, output_path)
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'Markdown to PDF conversion failed: {e}')

def convert_png_to_svg(file_path):
    try:
        output_path = file_path.replace('.png', '.svg')
        # Placeholder: Implement PNG to SVG vectorization
        with open(output_path, 'w') as file:
            file.write("<svg></svg>")
        return default_storage.url(output_path)
    except Exception as e:
        raise RuntimeError(f'PNG to SVG conversion failed: {e}')
